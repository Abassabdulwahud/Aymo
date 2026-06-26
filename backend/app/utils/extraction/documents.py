from pathlib import Path

import openpyxl
from docx import Document as WordDocument
from pptx import Presentation

from .base import ExtractionResult


def _extract_docx(file_path: Path) -> str:
    document = WordDocument(file_path)
    return "\n".join(paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip())


def _extract_pptx(file_path: Path) -> str:
    presentation = Presentation(file_path)
    slides = []
    for slide in presentation.slides:
        slide_parts = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "").strip()
            if text:
                slide_parts.append(text)
        if slide_parts:
            slides.append("\n".join(slide_parts))
    return "\n\n".join(slides)


def _extract_xlsx(file_path: Path) -> str:
    workbook = openpyxl.load_workbook(file_path, data_only=True)
    chunks = []
    for worksheet in workbook.worksheets:
        rows = []
        for row in worksheet.iter_rows(values_only=True):
            cells = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            chunks.append(f"[Sheet: {worksheet.title}]\n" + "\n".join(rows))
    return "\n\n".join(chunks)


def _extract_text_file(file_path: Path) -> str:
    return file_path.read_text(encoding="utf-8", errors="ignore").strip()


def extract_document_content(file_path: Path) -> ExtractionResult:
    suffix = file_path.suffix.lower()
    try:
        if suffix == ".docx":
            content = _extract_docx(file_path)
        elif suffix == ".pptx":
            content = _extract_pptx(file_path)
        elif suffix == ".xlsx":
            content = _extract_xlsx(file_path)
        else:
            content = _extract_text_file(file_path)
    except Exception as exc:
        return ExtractionResult(status="failed", content=None, error=f"Document extraction failed: {exc}")

    if not content:
        return ExtractionResult(status="failed", content=None, error="No text could be extracted from this document.")

    return ExtractionResult(status="completed", content=content)
